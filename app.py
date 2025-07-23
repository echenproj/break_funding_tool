import os
import re
from dateutil import parser

from flask import Flask, request, render_template
from flask import send_file
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from calculations import generate_cashflow_plot, compute_break_funding_cost
from extract_from_pdf import extract_loan_terms, chat_reply

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

FIELDS = [
    'effective_date', 'maturity_date', 'frequency', 'amortization_type',
    'loan_rate', 'balance'
]

def safe_field(field_name, extracted, form):
    """
    The safe_field() function is a helper that prioritizes user input, but falls back 
    to extracted values (like from a PDF), and if neither is available, returns a 
    blank string.
    """
    user_val = form.get(field_name)
    extracted_val = extracted.get(field_name, "")
    if user_val is not None and user_val.strip() != "":
        return user_val.strip()
    elif extracted_val not in (None, "", "None"):
        return extracted_val
    else:
        return ""

def normalize_date(date_str):
    try:
        dt = parser.parse(date_str)
        return dt.strftime("%Y-%m-%d")
    except Exception:
        return ""

def normalize_frequency(freq_str):
    freq_str = freq_str.strip().lower()
    mapping = {
        'monthly': 'monthly',
        'month': 'monthly',
        'quarterly': 'quarterly',
        'quarter': 'quarterly',
        '3m': 'quarterly',
        'semiannual': 'semiannual',
        'semi-annually': 'semiannual',
        '6m': 'semiannual',
        'annual': 'annual',
        'yearly': 'annual',
        '12m': 'annual'
    }
    return mapping.get(freq_str, '')

def normalize_amortization_type(amt_str):
    amt_str = amt_str.strip().lower()
    mapping = {
        'interest only': 'interest only',
        'io': 'interest only',
        'i/o': 'interest only',

        'equal': 'equal',
        'equal payment': 'equal',
        'annuity': 'equal',
        'level payment': 'equal',

        'linear': 'linear',
        'straight line': 'linear',
        'even principal': 'linear',
        'constant principal': 'linear',

        'custom': 'custom',
        'manual': 'custom',
        'user defined': 'custom',
    }
    return mapping.get(amt_str, '')


# Normalize/parse balance as float (handle commas, $ signs)
def parse_amount(s):
    if not s:
        return None
    s = re.sub(r'[^\d.]', '', s)
    try:
        return float(s)
    except ValueError:
        return None


@app.route('/', methods=['GET', 'POST'])
def index():
    extracted = {}
    extracted_quotes = {}
    data = {}
    plot_generated = False
    break_funding_cost = None
    error_message = None
    response_text=None

    if request.method == 'POST':
        action = request.form.get('action')  # Check whether it's 'upload' or 'calculate'

        if action == 'upload':
            # Only handle PDF upload and pre-fill the form
            pdf_file = request.files.get('pdf')
            if pdf_file and pdf_file.filename.endswith('.pdf'):
                extracted, extracted_quotes = extract_loan_terms(pdf_file)

            # Pre-fill fields with extracted values (if found), else blank
            for field in FIELDS:
                data[field] = extracted.get(field, '')

            # Extracted values from the PDF or user input
            effective_date = normalize_date(extracted.get('effective_date', ''))
            balance = parse_amount(extracted.get('balance', ''))

            # Set Prepayment Date default to Effective Date if missing
            if effective_date != "":
                prepayment_date = effective_date

            # Set Prepayment Amount default to half Balance if missing and Balance is valid
            if balance != "":
                prepayment_amount = float(balance / 2)

            # Leave prepayment fields empty
            data['prepayment_date'] = prepayment_date or ''
            data['prepayment_amount'] = prepayment_amount or ''

            data['effective_date'] = effective_date
            data['maturity_date'] = normalize_date(extracted.get('maturity_date', ''))
            data['frequency'] = normalize_frequency(safe_field('frequency', extracted, request.form))
            data['amortization_type'] = normalize_amortization_type(safe_field('amortization_type', extracted, request.form))

            return render_template('index.html', **data,
                                extracted_quotes=extracted_quotes,
                                plot_generated=False,
                                error_message=None,
                                break_funding_cost=None,
                                loading=False,
                                response_text=None)

        elif action == 'calculate':

            # Use safe_field to prioritize user input
            for field in FIELDS:
                data[field] = safe_field(field, extracted, request.form)

            data['prepayment_date'] = request.form.get('prepayment_date', '').strip()
            data['prepayment_amount'] = request.form.get('prepayment_amount', '').strip()

            data['effective_date'] = normalize_date(safe_field('effective_date', extracted, request.form))
            data['maturity_date'] = normalize_date(safe_field('maturity_date', extracted, request.form))
            data['frequency'] = normalize_frequency(safe_field('frequency', extracted, request.form))
            data['amortization_type'] = normalize_amortization_type(safe_field('amortization_type', extracted, request.form))

            required_fields = FIELDS + ['prepayment_date', 'prepayment_amount']
            if all(data.get(f) not in (None, '', 'None') for f in required_fields):
                try:
                    data['loan_rate'] = float(data['loan_rate'])
                    data['balance'] = float(data['balance'])
                    data['prepayment_amount'] = float(data['prepayment_amount'])
                except ValueError:
                    error_message = "Please enter valid numeric values."
                    return render_template("index.html", **data,
                                           extracted_quotes=extracted_quotes,
                                           plot_generated=False,
                                           error_message=error_message,
                                           break_funding_cost=None,
                                           loading=False,
                                           response_text=None)
                try:
                    generate_cashflow_plot(data)
                    break_funding_cost = compute_break_funding_cost(**data)
                    plot_generated = True
                except Exception as e:
                    error_message = f"Error generating plot: {e}"

            else:
                error_message = "Please fill in all required fields."

            return render_template('index.html', **data,
                                   extracted_quotes=extracted_quotes,
                                   plot_generated=plot_generated,
                                   error_message=error_message,
                                   break_funding_cost=break_funding_cost,
                                   loading=False,
                                   response_text=None)
        elif action == 'chat':
            for field in FIELDS:
                data[field] = request.form.get(field, '').strip()

            data['prepayment_date'] = request.form.get('prepayment_date', '').strip()
            data['prepayment_amount'] = request.form.get('prepayment_amount', '').strip()

            # Validate and normalize
            data['prepayment_date'] = request.form.get('prepayment_date', '').strip()
            data['prepayment_amount'] = request.form.get('prepayment_amount', '').strip()

            # Preserve previous outputs (pass them through the hidden fields)
            plot_generated = request.form.get('plot_generated') == 'True'
            try:
                break_funding_cost = float(request.form.get('break_funding_cost', 0))
            except (TypeError, ValueError):
                break_funding_cost = None

            # LLM response only
            user_input = request.form.get('user_input', '')
            response_text = None
            if user_input.strip():
                try:
                    response_text = chat_reply(user_input, break_funding_cost)
                except Exception as e:
                    error_message = f"LLM Error: {e}"

            return render_template('index.html', **data,
                                extracted_quotes=extracted_quotes,
                                plot_generated=plot_generated,
                                break_funding_cost=break_funding_cost,
                                error_message=error_message,
                                response_text=response_text,
                                loading=False)


    # GET request
    return render_template("index.html",
                           effective_date='',
                           maturity_date='',
                           frequency='',
                           amortization_type='',
                           loan_rate='',
                           balance='',
                           prepayment_date='',
                           prepayment_amount='',
                           extracted_quotes={},
                           plot_generated=False,
                           break_funding_cost=None,
                           error_message=None,
                           loading=False,
                           response_text=None)


@app.route('/download_ppt', methods=['POST'])
def download_ppt():
    # 1. Create presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # 2. Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "Break-Funding Analysis"
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    # 3. Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.5))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_run = subtitle_p.add_run()
    subtitle_run.text = "The following summarizes key loan details:"
    subtitle_run.font.size = Pt(20)
    subtitle_run.font.bold = True
    subtitle_run.font.color.rgb = RGBColor(0, 112, 192)  # Bank of America Blue

    # 4. Add loan details (bullets)
    loan_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(2))
    loan_tf = loan_textbox.text_frame
    loan_tf.clear()

    bullet_items = [
        ("Effective Date", request.form.get('effective_date', '')),
        ("Maturity Date", request.form.get('maturity_date', '')),
        ("Loan Rate", f"{request.form.get('loan_rate', '')}%"),
        ("Balance", f"${request.form.get('balance', '')}"),
    ]

    for label, value in bullet_items:
        p = loan_tf.add_paragraph()
        run = p.add_run()
        run.text = f"{label}: {value}"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.level = 0

    # 5. Break Funding Cost
    break_cost = request.form.get('break_funding_cost')
    if break_cost:
        cost_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(6), Inches(0.5))
        cost_tf = cost_box.text_frame
        cost_tf.clear()
        p = cost_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Break-Funding Cost: ${float(break_cost):,.2f}"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)

    # 6. Add chart image (bottom)
    img_path = "static/plot.png"
    img_height = Inches(3.5)
    img_width = Inches(7)
    img_top = prs.slide_height - img_height - Inches(0.3)
    img_left = (prs.slide_width - img_width) / 2
    slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

    # 7. Add logo (top-right)
    logo_path = "static/logo.png"
    logo_width = Inches(1.5)
    logo_top = Inches(0.2)
    logo_left = prs.slide_width - logo_width - Inches(0.3)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

    # 8. Save presentation to memory
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return send_file(
        ppt_io,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name='break_funding_analysis.pptx'
    )

if __name__ == '__main__':
    # port = int(os.environ.get("PORT", 5000))
    # app.run(debug=True, host='0.0.0.0', port=port)
    app.run(debug=False, host='0.0.0.0', port=int(os.environ.get("PORT", 10000)))